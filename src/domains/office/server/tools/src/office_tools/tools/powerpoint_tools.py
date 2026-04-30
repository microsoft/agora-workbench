"""
PowerPoint extraction tools.

All functions receive a ``file`` parameter that is a ``Path`` to a cached
PowerPoint file on disk (resolved by AssetResolutionMiddleware).
"""

from pathlib import Path
from typing import Any, Union

import olefile
from pptx import Presentation

_OLE2_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


def _check_irm(file: Path) -> None:
    """Raise ``ValueError`` if *file* is an OLE2 container (IRM-encrypted or legacy .ppt)."""
    with open(file, "rb") as f:
        header = f.read(8)
    if header != _OLE2_MAGIC:
        return
    if olefile.isOleFile(str(file)):
        with olefile.OleFileIO(str(file)) as ole:
            if ole.exists("\x09DRMContent"):
                raise ValueError(
                    f"File {file.name} is still IRM-encrypted. "
                    "IRM decryption may not have run — check that "
                    "the server has valid Azure RMS credentials."
                )
    raise ValueError(
        f"File {file.name} is a legacy OLE2 document (.ppt format). "
        "Only .pptx files are supported. Please re-save the file as .pptx."
    )


def extract_slides_text(file: Union[Path, Any]) -> dict:
    """
    Extract text from all slides in a PowerPoint presentation.

    Args:
        file: Path to the .pptx file.

    Returns:
        Dict with ``slides`` (list of dicts with slide_number and text)
        and ``slide_count``.
    """
    file = Path(file)
    _check_irm(file)
    prs = Presentation(str(file))

    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        slides.append(
            {
                "slide_number": idx,
                "text": "\n".join(texts),
            }
        )

    return {
        "slides": slides,
        "slide_count": len(slides),
    }


def extract_slide_notes(file: Union[Path, Any]) -> dict:
    """
    Extract speaker notes from all slides.

    Args:
        file: Path to the .pptx file.

    Returns:
        Dict with ``notes`` (list of dicts with slide_number and notes text)
        and ``slides_with_notes`` count.
    """
    file = Path(file)
    _check_irm(file)
    prs = Presentation(str(file))

    notes = []
    slides_with_notes = 0
    for idx, slide in enumerate(prs.slides, start=1):
        notes_text = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text.strip()

        if notes_text:
            slides_with_notes += 1

        notes.append(
            {
                "slide_number": idx,
                "notes": notes_text,
            }
        )

    return {
        "notes": notes,
        "slides_with_notes": slides_with_notes,
    }
